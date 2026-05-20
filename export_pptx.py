"""
export_pptx.py
Generates a branded Mews PowerPoint deck from dashboard data.
Requires:  python-pptx  matplotlib
Install:   pip install python-pptx matplotlib
"""

import io
import math
import matplotlib
matplotlib.use("Agg")          # headless – no display needed
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand colours ─────────────────────────────────────────────────────────────
BG        = RGBColor(0xFF, 0xFC, 0xF6)   # #FFFCF6  warm cream
BLACK     = RGBColor(0x26, 0x26, 0x26)   # #262626
PINK      = RGBColor(0xE8, 0x7D, 0xC2)   # #E87DC2
ORANGE    = RGBColor(0xFF, 0x6B, 0x00)   # #FF6B00
YELLOW    = RGBColor(0xD4, 0xE8, 0x33)   # #D4E833
WARM_GREY = RGBColor(0xE8, 0xE6, 0xDF)   # #E8E6DF
BLUE      = RGBColor(0xC8, 0xE5, 0xEE)   # #C8E5EE
MAUVE     = RGBColor(0xF0, 0xE0, 0xEE)   # #F0E0EE
DIVIDER   = RGBColor(0xE0, 0xDC, 0xD5)   # light rule

PALETTE_HEX = ["#E87DC2", "#FF6B00", "#D4E833", "#C8E5EE", "#F0E0EE", "#E8E6DF"]

# Slide size: widescreen 13.33 × 7.5 in
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Low-level helpers ─────────────────────────────────────────────────────────

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def _hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def _set_bg(slide, prs):
    """Fill slide background with brand cream."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_textbox(slide, text, left, top, width, height,
                 font_size=12, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def _add_rule(slide, top, color=DIVIDER):
    """Thin horizontal rule across the slide."""
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(0.4), top,
        Inches(12.93), top,
    )
    line.line.color.rgb = color
    line.line.width = Pt(0.5)


def _fig_to_image(fig):
    """Convert a matplotlib figure to a BytesIO PNG."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _add_image(slide, buf, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(buf, left, top, width, height)
    else:
        slide.shapes.add_picture(buf, left, top, width)


# ── Slide templates ───────────────────────────────────────────────────────────

def _cover_slide(prs, title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, prs)

    # Left accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.25), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PINK
    bar.line.fill.background()

    # Title
    _add_textbox(slide, title_text,
                 Inches(0.55), Inches(2.8), Inches(9), Inches(1.4),
                 font_size=36, bold=True, color=BLACK)

    if subtitle_text:
        _add_textbox(slide, subtitle_text,
                     Inches(0.55), Inches(4.0), Inches(9), Inches(0.6),
                     font_size=16, color=BLACK)

    # Date footer
    _add_textbox(slide, date.today().strftime("%B %d, %Y"),
                 Inches(0.55), Inches(6.8), Inches(5), Inches(0.4),
                 font_size=10, color=_rgb(0x88, 0x88, 0x88))
    return slide


def _section_slide(prs, section_title):
    """Coloured divider slide between sections."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

    _add_textbox(slide, section_title,
                 Inches(0.8), Inches(3.0), Inches(11), Inches(1.2),
                 font_size=32, bold=True, color=BG,
                 align=PP_ALIGN.LEFT)
    return slide


def _content_slide(prs, title):
    """Blank content slide with header area."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, prs)

    # Top accent strip
    strip = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    strip.fill.solid()
    strip.fill.fore_color.rgb = PINK
    strip.line.fill.background()

    # Slide title
    _add_textbox(slide, title,
                 Inches(0.4), Inches(0.12), Inches(12), Inches(0.5),
                 font_size=16, bold=True, color=BLACK)

    _add_rule(slide, Inches(0.72))
    return slide


# ── Matplotlib chart helpers ──────────────────────────────────────────────────
MEWS_BG = "#FFFCF6"


def _bar_chart(categories, values, color="#E87DC2", xlabel="", ylabel="",
               title="", value_fmt="{:.1f}%", figsize=(8, 4)):
    fig, ax = plt.subplots(figsize=figsize, facecolor=MEWS_BG)
    ax.set_facecolor(MEWS_BG)
    bars = ax.bar(categories, values, color=color, edgecolor="none", width=0.55)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2,
                bar.get_height() + max(values) * 0.01,
                value_fmt.format(val), ha="center", va="bottom",
                fontsize=8, color="#262626")
    ax.set_xlabel(xlabel, fontsize=9)
    ax.set_ylabel(ylabel, fontsize=9)
    ax.set_title(title, fontsize=11, fontweight="bold", color="#262626", pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#CCCCCC")
    ax.tick_params(colors="#262626", labelsize=8)
    ax.yaxis.grid(True, color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig


def _hbar_chart(categories, values, color="#E87DC2", xlabel="", title="",
                value_fmt="{:.1f}%", figsize=(7, 4)):
    fig, ax = plt.subplots(figsize=figsize, facecolor=MEWS_BG)
    ax.set_facecolor(MEWS_BG)
    y = range(len(categories))
    bars = ax.barh(list(y), values, color=color, edgecolor="none", height=0.55)
    ax.set_yticks(list(y))
    ax.set_yticklabels(categories, fontsize=8)
    for bar, val in zip(bars, values):
        ax.text(val + max(values) * 0.01, bar.get_y() + bar.get_height() / 2,
                value_fmt.format(val), va="center", fontsize=8, color="#262626")
    ax.set_xlabel(xlabel, fontsize=9)
    ax.set_title(title, fontsize=11, fontweight="bold", color="#262626", pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#CCCCCC")
    ax.tick_params(colors="#262626", labelsize=8)
    ax.xaxis.grid(True, color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return fig


def _pie_chart(labels, values, colors=None, title="", figsize=(5, 4)):
    if colors is None:
        colors = PALETTE_HEX[:len(labels)]
    fig, ax = plt.subplots(figsize=figsize, facecolor=MEWS_BG)
    ax.set_facecolor(MEWS_BG)
    wedges, texts, autotexts = ax.pie(
        values, labels=None, colors=colors,
        autopct="%1.1f%%", startangle=90,
        wedgeprops=dict(edgecolor="white", linewidth=1.5))
    for at in autotexts:
        at.set_fontsize(8)
        at.set_color("#262626")
    ax.legend(wedges, labels, loc="center left",
              bbox_to_anchor=(1, 0.5), fontsize=8,
              frameon=False)
    ax.set_title(title, fontsize=11, fontweight="bold", color="#262626", pad=8)
    fig.tight_layout()
    return fig


def _line_chart(df_pivot, ylabel="", title="", colors=None, figsize=(9, 4)):
    """df_pivot: index=x-axis, columns=series."""
    if colors is None:
        colors = PALETTE_HEX
    fig, ax = plt.subplots(figsize=figsize, facecolor=MEWS_BG)
    ax.set_facecolor(MEWS_BG)
    for i, col in enumerate(df_pivot.columns):
        lw = 2.5 if str(col) == "Global" else 1.5
        ls = "--" if str(col) == "Global" else "-"
        ax.plot(df_pivot.index, df_pivot[col],
                label=str(col), color=colors[i % len(colors)],
                linewidth=lw, linestyle=ls)
    ax.set_ylabel(ylabel, fontsize=9)
    ax.set_title(title, fontsize=11, fontweight="bold", color="#262626", pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#CCCCCC")
    ax.tick_params(colors="#262626", labelsize=7, rotation=30)
    ax.yaxis.grid(True, color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.legend(fontsize=8, frameon=False, loc="upper left",
              bbox_to_anchor=(1.01, 1))
    fig.tight_layout()
    return fig


def _grouped_bar(categories, series_dict, title="", ylabel="",
                 value_fmt="{:.1f}", figsize=(9, 4)):
    """series_dict = {label: [values...]}"""
    n_groups = len(categories)
    n_series = len(series_dict)
    x = np.arange(n_groups)
    width = 0.7 / n_series
    fig, ax = plt.subplots(figsize=figsize, facecolor=MEWS_BG)
    ax.set_facecolor(MEWS_BG)
    for i, (label, vals) in enumerate(series_dict.items()):
        offset = (i - n_series / 2 + 0.5) * width
        bars = ax.bar(x + offset, vals, width,
                      label=label,
                      color=PALETTE_HEX[i % len(PALETTE_HEX)],
                      edgecolor="none")
    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=8, rotation=15, ha="right")
    ax.set_ylabel(ylabel, fontsize=9)
    ax.set_title(title, fontsize=11, fontweight="bold", color="#262626", pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#CCCCCC")
    ax.tick_params(colors="#262626", labelsize=8)
    ax.yaxis.grid(True, color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.legend(fontsize=8, frameon=False)
    fig.tight_layout()
    return fig


def _stacked_bar_pct(df_monthly, title="", figsize=(10, 4)):
    """Stacked 100% bar for P&L monthly."""
    months = df_monthly["month"].dt.strftime("%b %Y").unique()
    categories = df_monthly["category"].unique()
    totals = df_monthly.groupby("month")["net_revenue_eur"].transform("sum")
    df_p = df_monthly.copy()
    df_p["pct"] = df_p["net_revenue_eur"] / totals * 100

    fig, ax = plt.subplots(figsize=figsize, facecolor=MEWS_BG)
    ax.set_facecolor(MEWS_BG)
    bottom = np.zeros(len(months))
    month_labels = sorted(df_p["month"].unique())

    for i, cat in enumerate(categories):
        sub = df_p[df_p["category"] == cat].sort_values("month")
        vals = [sub[sub["month"] == m]["pct"].values[0]
                if m in sub["month"].values else 0
                for m in month_labels]
        ax.bar(range(len(month_labels)), vals, bottom=bottom,
               label=cat, color=PALETTE_HEX[i % len(PALETTE_HEX)],
               edgecolor="none")
        bottom += np.array(vals)

    ax.set_xticks(range(len(month_labels)))
    ax.set_xticklabels([m.strftime("%b %Y") for m in month_labels],
                       fontsize=7, rotation=45, ha="right")
    ax.set_ylabel("% of Revenue", fontsize=9)
    ax.set_title(title, fontsize=11, fontweight="bold", color="#262626", pad=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.spines[["left", "bottom"]].set_color("#CCCCCC")
    ax.tick_params(colors="#262626", labelsize=7)
    ax.legend(fontsize=7, frameon=False, loc="upper left",
              bbox_to_anchor=(1.01, 1))
    fig.tight_layout()
    return fig


# ── KPI tile helper (rendered as a table-like shape) ─────────────────────────

def _kpi_tiles_slide(prs, slide_title, kpi_rows):
    """
    kpi_rows: list of (label, value_str, delta_str_or_None)
    Renders up to 5 KPI boxes in a row.
    """
    slide = _content_slide(prs, slide_title)
    n = len(kpi_rows)
    box_w = Inches(12.0 / n)
    box_h = Inches(2.0)
    top   = Inches(2.0)

    for i, (label, value, delta) in enumerate(kpi_rows):
        left = Inches(0.4) + i * box_w

        # Box background
        rect = slide.shapes.add_shape(1, left, top, box_w - Inches(0.15), box_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = WARM_GREY
        rect.line.fill.background()

        # Label
        _add_textbox(slide, label,
                     left + Inches(0.12), top + Inches(0.12),
                     box_w - Inches(0.27), Inches(0.4),
                     font_size=9, color=BLACK)
        # Value
        _add_textbox(slide, value,
                     left + Inches(0.12), top + Inches(0.5),
                     box_w - Inches(0.27), Inches(0.7),
                     font_size=22, bold=True, color=BLACK)
        # Delta
        if delta:
            is_pos = delta.startswith("+")
            d_color = _rgb(0x2E, 0x7D, 0x32) if is_pos else _rgb(0xC6, 0x28, 0x28)
            _add_textbox(slide, delta,
                         left + Inches(0.12), top + Inches(1.25),
                         box_w - Inches(0.27), Inches(0.4),
                         font_size=9, color=d_color)
    return slide


def _annual_tiles_slide(prs, slide_title, df_ann, dim_col, items):
    """
    Annual metric tiles: occupancy / ADR / RevPAR × year columns.
    items = list of dimension values to show (plus Global).
    """
    slide = _content_slide(prs, slide_title)
    show_items = ["Global"] + items

    metrics = [
        ("occupancy", "Occupancy (%)", lambda v: f"{v:.1f}%"),
        ("adr",       "Avg ADR (€)",   lambda v: f"€{v:,.0f}"),
        ("revpar",    "Avg RevPAR (€)",lambda v: f"€{v:,.0f}"),
    ]

    row_h = Inches(1.55)
    col_labels = ["Segment / Class", "2024", "2025", "2026 (YTD)"]
    col_widths = [Inches(2.8), Inches(2.8), Inches(2.8), Inches(2.8)]
    start_top = Inches(0.85)
    start_left = Inches(0.4)

    for mi, (metric, mlabel, fmt) in enumerate(metrics):
        base_top = start_top + mi * (row_h + Inches(0.08))

        # Metric label on left
        _add_textbox(slide, mlabel,
                     start_left, base_top,
                     Inches(12.5), Inches(0.32),
                     font_size=10, bold=True, color=BLACK)

        # Header row
        for ci, (cl, cw) in enumerate(zip(col_labels, col_widths)):
            lft = start_left + sum(col_widths[:ci])
            _add_textbox(slide, cl,
                         lft, base_top + Inches(0.32),
                         cw - Inches(0.05), Inches(0.28),
                         font_size=8, bold=True, color=BLACK)

        # Data rows
        for ri, item in enumerate(show_items):
            row_df = df_ann[df_ann[dim_col] == item]
            item_top = base_top + Inches(0.32) + Inches(0.28) + ri * Inches(0.32)

            # Row bg
            row_bg = slide.shapes.add_shape(
                1, start_left, item_top,
                sum(col_widths) - Inches(0.05), Inches(0.28))
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = (WARM_GREY if ri % 2 == 0
                                          else RGBColor(0xFF, 0xFC, 0xF6))
            row_bg.line.fill.background()

            # Item name
            _add_textbox(slide, "🌍 Global" if item == "Global" else f"  {item}",
                         start_left, item_top,
                         col_widths[0] - Inches(0.05), Inches(0.28),
                         font_size=8, bold=(item == "Global"), color=BLACK)

            # Year columns
            for ci, yr in enumerate([2024, 2025, 2026], 1):
                yr_row = row_df[row_df["year"] == yr]
                lft = start_left + sum(col_widths[:ci])
                if yr_row.empty or pd.isna(yr_row[metric].iloc[0]):
                    val_str = "—"
                else:
                    val = float(yr_row[metric].iloc[0])
                    prev = row_df[row_df["year"] == yr - 1]
                    if not prev.empty and not pd.isna(prev[metric].iloc[0]):
                        pv = float(prev[metric].iloc[0])
                        if metric == "occupancy":
                            delta_str = f" ({val-pv:+.1f}pp)"
                        else:
                            delta_str = f" ({(val-pv)/pv*100:+.1f}%)"
                    else:
                        delta_str = ""
                    val_str = fmt(val) + delta_str
                _add_textbox(slide, val_str,
                             lft, item_top,
                             col_widths[ci] - Inches(0.05), Inches(0.28),
                             font_size=8, color=BLACK)
    return slide


# ── Main build function ───────────────────────────────────────────────────────

def build_pptx(filters: dict,
               kpis: dict,
               ytd: dict,
               df_trends: pd.DataFrame,
               df_otb: pd.DataFrame,
               otb_g: dict,
               df_reg_ann: pd.DataFrame,
               df_reg_mon: pd.DataFrame,
               selected_territory: list,
               df_hc_ann: pd.DataFrame,
               df_hc_mon: pd.DataFrame,
               selected_classes: list,
               beh: dict,
               df_los: pd.DataFrame,
               df_gs: pd.DataFrame,
               df_lt: pd.DataFrame,
               df_ch: pd.DataFrame,
               df_pay: pd.DataFrame,
               df_card: pd.DataFrame,
               df_canc: pd.DataFrame,
               df_pnl_monthly: pd.DataFrame,
               df_pnl_mix: pd.DataFrame,
               df_sqm: pd.DataFrame,
               df_ci_dow: pd.DataFrame,
               df_co_dow: pd.DataFrame,
               df_pre: pd.DataFrame,
               df_hotel: pd.DataFrame,
               df_uch: pd.DataFrame,
               df_uval: pd.DataFrame,
               df_ucat: pd.DataFrame,
               ) -> bytes:

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Filter description for subtitle
    terr_str = ", ".join(filters.get("territory") or ["Global"])
    seg_str  = ", ".join(filters.get("segment")   or ["All segments"])
    subtitle = (f"Territory: {terr_str}  |  Segment: {seg_str}  |  "
                f"{filters['start']} → {filters['end']}")

    # ── Cover ──────────────────────────────────────────────────────────────────
    _cover_slide(prs, "Hotel Performance Report", subtitle)

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 1 — Market KPIs
    # ═══════════════════════════════════════════════════════════════════════════
    _section_slide(prs, "Market KPIs")

    # --- Slide: KPI tiles + YTD growth ----------------------------------------
    kpi_rows = [
        ("Enterprises",        f"{kpis.get('enterprises', 0):,}",       None),
        ("Avg ADR (€)",        f"€{kpis.get('adr', 0):,.2f}",           None),
        ("Avg Occupancy",      f"{kpis.get('occupancy', 0):.1f}%",      None),
        ("Avg RevPAR (€)",     f"€{kpis.get('revpar', 0):,.2f}",        None),
        ("Total Reservations", f"{kpis.get('reservations', 0):,}",      None),
    ]
    _kpi_tiles_slide(prs, "Key Metrics", kpi_rows)

    # --- Slide: YTD Growth tiles -----------------------------------------------
    if ytd:
        ytd_rows = [
            ("Occupancy 2026 YTD", f"{ytd['occ_2026']:.1f}%",
             f"{ytd['occ_chg']:+.1f}% vs 2025 ({ytd['occ_2025']:.1f}%)"),
            ("ADR 2026 YTD", f"€{ytd['adr_2026']:,.2f}",
             f"{ytd['adr_chg']:+.1f}% vs 2025"),
            ("RevPAR 2026 YTD", f"€{ytd['revpar_2026']:,.2f}",
             f"{ytd['revpar_chg']:+.1f}% vs 2025"),
        ]
        _kpi_tiles_slide(prs, f"YTD Growth — 2026 vs 2025", ytd_rows)

    # --- Slide: Historical trends (3 line charts side by side) ----------------
    if not df_trends.empty:
        slide = _content_slide(prs, "Historical Performance (7-day rolling average)")
        df_trends["year"] = df_trends["year"].astype(str)
        years = sorted(df_trends["year"].unique())
        year_colors = {y: PALETTE_HEX[i % len(PALETTE_HEX)]
                       for i, y in enumerate(years)}
        chart_top  = Inches(0.85)
        chart_h    = Inches(5.8)
        chart_w    = Inches(4.1)
        for ci, (metric, label) in enumerate([
            ("occupancy", "Occupancy (%)"),
            ("adr",       "ADR (€)"),
            ("revpar",    "RevPAR (€)"),
        ]):
            pivot = df_trends.pivot_table(
                index="day_of_year", columns="year", values=metric, aggfunc="mean")
            fig = _line_chart(pivot, ylabel=label, title=label,
                              colors=[year_colors[y] for y in pivot.columns],
                              figsize=(5.5, 4))
            buf = _fig_to_image(fig)
            _add_image(slide, buf,
                       Inches(0.3) + ci * chart_w, chart_top,
                       chart_w - Inches(0.1), chart_h)

    # --- Slide: OTB trends ----------------------------------------------------
    if not df_otb.empty:
        slide = _content_slide(prs, "On The Books — Next 9 Months vs Last Year")
        otb_colors = {"This Year (OTB)": "#E87DC2", "Last Year (OTB)": "#E8E6DF"}
        chart_w = Inches(4.1)
        for ci, (metric, label) in enumerate([
            ("occupancy", "Occupancy (%)"),
            ("adr",       "ADR (€)"),
            ("revpar",    "RevPAR (€)"),
        ]):
            pivot = df_otb.pivot_table(
                index="days_ahead", columns="year_label", values=metric,
                aggfunc="mean")
            colors = [otb_colors.get(c, PALETTE_HEX[i])
                      for i, c in enumerate(pivot.columns)]
            fig = _line_chart(pivot, ylabel=label, title=label,
                              colors=colors, figsize=(5.5, 4))
            buf = _fig_to_image(fig)
            _add_image(slide, buf,
                       Inches(0.3) + ci * chart_w, Inches(0.85),
                       chart_w - Inches(0.1), Inches(5.8))

    # --- Slide: OTB Growth tiles -----------------------------------------------
    if otb_g:
        otb_rows = [
            ("OTB Occupancy 2026", f"{otb_g['occ_ty']:.1f}%",
             f"{otb_g['occ_chg']:+.1f}% vs 2025 ({otb_g['occ_ly']:.1f}%)"),
            ("OTB ADR 2026", f"€{otb_g['adr_ty']:,.2f}",
             f"{otb_g['adr_chg']:+.1f}% vs 2025"),
            ("OTB RevPAR 2026", f"€{otb_g['revpar_ty']:,.2f}",
             f"{otb_g['revpar_chg']:+.1f}% vs 2025"),
        ]
        snap = otb_g.get("snap_date", "latest snapshot")
        _kpi_tiles_slide(prs, f"OTB Growth — next 9 months from {snap}", otb_rows)

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 2 — Regional Overview
    # ═══════════════════════════════════════════════════════════════════════════
    _section_slide(prs, "Regional Overview")

    if not df_reg_ann.empty:
        _annual_tiles_slide(prs, "Annual Averages by Territory",
                            df_reg_ann, "territory", selected_territory)

    if not df_reg_mon.empty:
        terr_colors = {t: PALETTE_HEX[i % len(PALETTE_HEX)]
                       for i, t in enumerate(selected_territory)}
        terr_colors["Global"] = "#262626"
        show_t = ["Global"] + selected_territory
        df_plot = df_reg_mon[df_reg_mon["territory"].isin(show_t)].copy()

        for metric, label in [
            ("occupancy", "Occupancy (%)"),
            ("adr",       "ADR (€)"),
            ("revpar",    "RevPAR (€)"),
        ]:
            slide = _content_slide(prs, f"Monthly Trends by Territory — {label}")
            if not df_plot.empty:
                pivot = df_plot.pivot_table(
                    index="month", columns="territory", values=metric)
                colors = [terr_colors.get(c, PALETTE_HEX[i])
                          for i, c in enumerate(pivot.columns)]
                fig = _line_chart(pivot, ylabel=label, title=label,
                                  colors=colors, figsize=(11, 5))
                buf = _fig_to_image(fig)
                _add_image(slide, buf,
                           Inches(0.4), Inches(0.85),
                           Inches(12.0), Inches(5.9))

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 3 — Hotel Class
    # ═══════════════════════════════════════════════════════════════════════════
    _section_slide(prs, "Hotel Class")

    if not df_hc_ann.empty:
        _annual_tiles_slide(prs, "Annual Averages by Hotel Class",
                            df_hc_ann, "hotel_class", selected_classes)

    if not df_hc_mon.empty:
        hc_colors = {c: PALETTE_HEX[i % len(PALETTE_HEX)]
                     for i, c in enumerate(selected_classes)}
        hc_colors["Global"] = "#262626"
        show_hc = ["Global"] + selected_classes
        df_hc_plot = df_hc_mon[df_hc_mon["hotel_class"].isin(show_hc)].copy()

        for metric, label in [
            ("occupancy", "Occupancy (%)"),
            ("adr",       "ADR (€)"),
            ("revpar",    "RevPAR (€)"),
        ]:
            slide = _content_slide(prs, f"Monthly Trends by Hotel Class — {label}")
            if not df_hc_plot.empty:
                pivot = df_hc_plot.pivot_table(
                    index="month", columns="hotel_class", values=metric)
                colors = [hc_colors.get(c, PALETTE_HEX[i])
                          for i, c in enumerate(pivot.columns)]
                fig = _line_chart(pivot, ylabel=label, title=label,
                                  colors=colors, figsize=(11, 5))
                buf = _fig_to_image(fig)
                _add_image(slide, buf,
                           Inches(0.4), Inches(0.85),
                           Inches(12.0), Inches(5.9))

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 4 — Booking Behaviour
    # ═══════════════════════════════════════════════════════════════════════════
    _section_slide(prs, "Booking Behaviour")

    # --- Length of stay --------------------------------------------------------
    if not df_los.empty:
        slide = _content_slide(prs, "Reservations by Length of Stay")
        df_los["pct"] = (df_los["reservations"] /
                         df_los["reservations"].sum() * 100)
        fig = _bar_chart(df_los["los_bucket"].astype(str).tolist(),
                         df_los["pct"].tolist(),
                         color="#FF6B00", ylabel="% of Reservations",
                         title="Length of Stay", figsize=(9, 4.5))
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(1.5), Inches(0.85),
                   Inches(10.0), Inches(5.9))

    # --- Group size ------------------------------------------------------------
    if not df_gs.empty:
        slide = _content_slide(prs, "Reservations by Group Size")
        df_gs["pct"] = df_gs["reservations"] / df_gs["reservations"].sum() * 100
        fig = _bar_chart(df_gs["group_size_bucket"].astype(str).tolist(),
                         df_gs["pct"].tolist(),
                         color="#D4E833", ylabel="% of Reservations",
                         title="Group Size", figsize=(9, 4.5))
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(1.5), Inches(0.85),
                   Inches(10.0), Inches(5.9))

    # --- Lead time -------------------------------------------------------------
    if not df_lt.empty:
        slide = _content_slide(prs, "Reservations by Lead Time")
        df_lt["pct"] = df_lt["reservations"] / df_lt["reservations"].sum() * 100
        fig = _bar_chart(df_lt["lead_time_bucket"].astype(str).tolist(),
                         df_lt["pct"].tolist(),
                         color="#E87DC2", ylabel="% of Reservations",
                         title="Lead Time", figsize=(9, 4.5))
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(1.5), Inches(0.85),
                   Inches(10.0), Inches(5.9))

    # --- Booking channel -------------------------------------------------------
    if not df_ch.empty:
        slide = _content_slide(prs, "Reservations by Booking Channel")
        fig = _pie_chart(df_ch["channel"].tolist(),
                         df_ch["reservations"].tolist(),
                         title="Booking Channel")
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(3.0), Inches(0.85),
                   Inches(7.0), Inches(5.9))

    # --- Payment type ----------------------------------------------------------
    if not df_pay.empty:
        slide = _content_slide(prs, "Payment Type Breakdown")
        fig = _pie_chart(df_pay["category"].tolist(),
                         df_pay["transactions"].tolist(),
                         title="Payment Type")
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(3.0), Inches(0.85),
                   Inches(7.0), Inches(5.9))

    # --- Card network ----------------------------------------------------------
    if not df_card.empty:
        slide = _content_slide(prs, "Card Network Breakdown")
        card_colors = {"Visa": "#C8E5EE", "Mastercard": "#E87DC2",
                       "Amex": "#D4E833", "Other": "#E8E6DF"}
        colors = [card_colors.get(n, "#E87DC2")
                  for n in df_card["card_network"].tolist()]
        fig = _pie_chart(df_card["card_network"].tolist(),
                         df_card["transactions"].tolist(),
                         colors=colors, title="Card Network")
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(3.0), Inches(0.85),
                   Inches(7.0), Inches(5.9))

    # --- Check-in day of week -------------------------------------------------
    if not df_ci_dow.empty and not df_co_dow.empty:
        slide = _content_slide(prs, "Arrivals & Departures by Day of Week")
        fig_ci = _bar_chart(df_ci_dow["day_of_week"].tolist(),
                            df_ci_dow["pct"].tolist(),
                            color="#E87DC2", ylabel="% of Arrivals",
                            title="Check-in Day", figsize=(6, 4))
        fig_co = _bar_chart(df_co_dow["day_of_week"].tolist(),
                            df_co_dow["pct"].tolist(),
                            color="#C8E5EE", ylabel="% of Departures",
                            title="Check-out Day", figsize=(6, 4))
        _add_image(slide, _fig_to_image(fig_ci),
                   Inches(0.3), Inches(0.85), Inches(6.3), Inches(5.9))
        _add_image(slide, _fig_to_image(fig_co),
                   Inches(6.7), Inches(0.85), Inches(6.3), Inches(5.9))

    # --- Cancellations --------------------------------------------------------
    if not df_canc.empty:
        slide = _content_slide(prs, "Cancellation Insights")
        years   = df_canc["year"].astype(str).tolist()
        cr_vals = df_canc["cancel_rate"].tolist()
        cw_vals = df_canc["avg_cancel_window"].fillna(0).tolist()

        fig = _grouped_bar(years,
                           {"Cancellation Rate (%)": cr_vals},
                           title="Cancellation Rate (%)",
                           ylabel="%", figsize=(6, 4))
        fig2 = _grouped_bar(years,
                            {"Avg Days Before Arrival": cw_vals},
                            title="Avg Cancellation Window (days)",
                            ylabel="days", figsize=(6, 4))
        _add_image(slide, _fig_to_image(fig),
                   Inches(0.3), Inches(0.85), Inches(6.3), Inches(5.9))
        _add_image(slide, _fig_to_image(fig2),
                   Inches(6.7), Inches(0.85), Inches(6.3), Inches(5.9))

    # ═══════════════════════════════════════════════════════════════════════════
    # SECTION 5 — Hotel Performance
    # ═══════════════════════════════════════════════════════════════════════════
    _section_slide(prs, "Hotel Performance")

    # --- P&L monthly stacked bar ----------------------------------------------
    if not df_pnl_monthly.empty:
        slide = _content_slide(prs, "Monthly Revenue by Department (%)")
        fig = _stacked_bar_pct(df_pnl_monthly,
                               title="Monthly Revenue by Department (%)",
                               figsize=(11, 5))
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(0.4), Inches(0.85),
                   Inches(12.0), Inches(5.9))

    # --- Revenue mix pie -------------------------------------------------------
    if not df_pnl_mix.empty:
        slide = _content_slide(prs, "Revenue Mix by Department")
        df_pie = df_pnl_mix[df_pnl_mix["category"] != "Not Assigned"].copy()
        fig = _pie_chart(df_pie["category"].tolist(),
                         df_pie["net_revenue_eur"].tolist(),
                         title="Revenue Mix")
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(2.5), Inches(0.85),
                   Inches(8.0), Inches(5.9))

    # --- Revenue per sqm -------------------------------------------------------
    if not df_sqm.empty:
        slide = _content_slide(prs, "Revenue per m² by Department")
        sorted_df = df_sqm.sort_values("revenue_per_sqm_eur")
        fig = _hbar_chart(
            sorted_df["category"].tolist(),
            sorted_df["revenue_per_sqm_eur"].tolist(),
            color="#E87DC2", xlabel="€ per m²",
            title="Avg Revenue per m²",
            value_fmt="€{:.2f}", figsize=(9, 4.5))
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(1.5), Inches(0.85),
                   Inches(10.0), Inches(5.9))

    # --- Check-in methods -----------------------------------------------------
    if not df_hotel.empty:
        slide = _content_slide(prs, "At-Hotel Check-in Methods")
        fig = _pie_chart(df_hotel["checkin_method"].tolist(),
                         df_hotel["reservations"].tolist(),
                         title="Check-in Method (at hotel)")
        buf = _fig_to_image(fig)
        _add_image(slide, buf, Inches(3.0), Inches(0.85),
                   Inches(7.0), Inches(5.9))

    # --- Upsells by channel ---------------------------------------------------
    if not df_uch.empty:
        slide = _content_slide(prs, "Upsell Performance")
        fig_left = _pie_chart(df_uch["channel"].tolist(),
                              df_uch["upsells"].tolist(),
                              title="Upsells by Channel")
        buf_left = _fig_to_image(fig_left)

        if not df_uval.empty:
            fig_right = _bar_chart(
                df_uval["channel"].tolist(),
                df_uval["avg_value"].tolist(),
                color="#FF6B00",
                ylabel="Avg Value (€)",
                title="Avg Upsell Value / Reservation (€)",
                value_fmt="€{:.2f}", figsize=(6, 4))
            buf_right = _fig_to_image(fig_right)
            _add_image(slide, buf_left,
                       Inches(0.3), Inches(0.85), Inches(6.3), Inches(5.9))
            _add_image(slide, buf_right,
                       Inches(6.7), Inches(0.85), Inches(6.3), Inches(5.9))
        else:
            _add_image(slide, buf_left,
                       Inches(3.0), Inches(0.85), Inches(7.0), Inches(5.9))

    # --- Top upsell categories ------------------------------------------------
    if not df_ucat.empty:
        slide = _content_slide(prs, "Top Upsell Categories")
        df_ucat_s = df_ucat.sort_values("upsells", ascending=True)
        df_ucat_s["pct"] = (df_ucat_s["upsells"] /
                            df_ucat_s["upsells"].sum() * 100)
        fig_l = _hbar_chart(
            df_ucat_s["category"].tolist(),
            df_ucat_s["pct"].tolist(),
            color="#C8E5EE", xlabel="% of Upsells",
            title="Top Categories by Volume", value_fmt="{:.1f}%",
            figsize=(6, 4))
        fig_r = _hbar_chart(
            df_ucat_s["category"].tolist(),
            df_ucat_s["avg_value_eur"].tolist(),
            color="#F0E0EE", xlabel="Avg Value (€)",
            title="Avg Value per Upsell (€)", value_fmt="€{:.2f}",
            figsize=(6, 4))
        _add_image(slide, _fig_to_image(fig_l),
                   Inches(0.3), Inches(0.85), Inches(6.3), Inches(5.9))
        _add_image(slide, _fig_to_image(fig_r),
                   Inches(6.7), Inches(0.85), Inches(6.3), Inches(5.9))

    # ── Serialise to bytes ────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()